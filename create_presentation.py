from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Design & Color Palette ---
# A modern corporate color scheme
COLOR_DARK_BLUE = RGBColor(0, 45, 86)      # #002D56
COLOR_LIGHT_BLUE = RGBColor(0, 123, 255)   # #007BFF
COLOR_GOLD = RGBColor(255, 193, 7)         # #FFC107
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GREY = RGBColor(108, 117, 125)      # #6C757D
COLOR_LIGHT_GREY = RGBColor(248, 249, 250) # #F8F9FA
FONT_FAMILY = "Arial"

def set_slide_background(slide, color):
    """Sets a solid color background for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_footer(slide, text):
    """Adds a consistent footer to each slide."""
    footer_shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.0), Inches(9), Inches(0.4)
    )
    footer_para = footer_shape.text_frame.paragraphs[0]
    footer_para.text = text
    footer_para.font.name = FONT_FAMILY
    footer_para.font.size = Pt(10)
    footer_para.font.color.rgb = COLOR_GREY
    footer_para.alignment = PP_ALIGN.LEFT

# --- Presentation Generation ---
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# --- Slide 1: Title Slide ---
slide1_layout = prs.slide_layouts[5] # Blank layout
slide1 = prs.slides.add_slide(slide1_layout)
set_slide_background(slide1, COLOR_DARK_BLUE)

# Title
title_shape = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
title_frame = title_shape.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "🎯 MODERN HR LEAVE MANAGEMENT SYSTEM"
title_para.font.name = FONT_FAMILY
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = COLOR_WHITE
title_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_shape = slide1.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(1.0))
subtitle_para = subtitle_shape.text_frame.paragraphs[0]
subtitle_para.text = "Al Ahlia Contracting Group"
subtitle_para.font.name = FONT_FAMILY
subtitle_para.font.size = Pt(28)
subtitle_para.font.color.rgb = COLOR_GOLD
subtitle_para.alignment = PP_ALIGN.CENTER

# Key points
points_text = "• Kuwait Labor Law Compliant\n• Real-time Analytics Dashboard\n• 94% Faster Processing\n• 98% Employee Satisfaction"
points_shape = slide1.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(2.0))
points_frame = points_shape.text_frame
points_frame.word_wrap = True
for line in points_text.split('\n'):
    p = points_frame.add_paragraph()
    p.text = line
    p.font.name = FONT_FAMILY
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.5

# --- Slide 2: Holiday Statistics & Analytics ---
slide2_layout = prs.slide_layouts[5]
slide2 = prs.slides.add_slide(slide2_layout)
set_slide_background(slide2, COLOR_LIGHT_GREY)

# Title
title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
p2 = title2.text_frame.paragraphs[0]
p2.text = "📊 HOLIDAY INTELLIGENCE DASHBOARD"
p2.font.name = FONT_FAMILY
p2.font.bold = True
p2.font.size = Pt(32)
p2.font.color.rgb = COLOR_DARK_BLUE

# Leave Distribution
dist_title = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(0.5))
dist_title.text_frame.paragraphs[0].text = "📈 LEAVE DISTRIBUTION 2024"
dist_title.text_frame.paragraphs[0].font.bold = True
dist_title.text_frame.paragraphs[0].font.size = Pt(18)
dist_title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_BLUE

dist_data = {"Q1": 35, "Q2": 45, "Q3": 20, "Q4": 15}
y_pos = 2.0
for q, val in dist_data.items():
    slide2.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(1), Inches(0.4)).text_frame.paragraphs[0].text = f"{q}:"
    bar_bg = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(y_pos + 0.1), Inches(2.5), Inches(0.2))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = RGBColor(200, 200, 200)
    bar_bg.line.fill.background()
    bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(y_pos + 0.1), Inches(2.5 * (val/100)), Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    bar.line.fill.background()
    slide2.shapes.add_textbox(Inches(4.1), Inches(y_pos), Inches(1), Inches(0.4)).text_frame.paragraphs[0].text = f"{val}%"
    y_pos += 0.5

# Department Utilization
dept_title = slide2.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(0.5))
dept_title.text_frame.paragraphs[0].text = "🏢 DEPARTMENT UTILIZATION"
dept_title.text_frame.paragraphs[0].font.bold = True
dept_title.text_frame.paragraphs[0].font.size = Pt(18)
dept_title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_BLUE

dept_data = {"IT": 42, "HR": 35, "Sales": 48, "Finance": 28}
max_days = max(dept_data.values())
y_pos = 2.0
for dept, val in dept_data.items():
    slide2.shapes.add_textbox(Inches(5.5), Inches(y_pos), Inches(1), Inches(0.4)).text_frame.paragraphs[0].text = f"{dept}:"
    bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(y_pos + 0.1), Inches(2.5 * (val/max_days)), Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_GOLD
    bar.line.fill.background()
    slide2.shapes.add_textbox(Inches(9.1), Inches(y_pos), Inches(1), Inches(0.4)).text_frame.paragraphs[0].text = f"{val} days"
    y_pos += 0.5
    
# Business Impact
impact_title = slide2.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.5))
impact_title.text_frame.paragraphs[0].text = "🎯 BUSINESS IMPACT"
impact_title.text_frame.paragraphs[0].font.bold = True
impact_title.text_frame.paragraphs[0].font.size = Pt(18)
impact_title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_BLUE

impact_text = "• 40% fewer staffing conflicts\n• 25% better resource planning\n• 60% faster decisions"
impact_shape = slide2.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1.5))
impact_frame = impact_shape.text_frame
impact_frame.word_wrap = True
for line in impact_text.split('\n'):
    p = impact_frame.add_paragraph()
    p.text = line
    p.font.name = FONT_FAMILY
    p.font.size = Pt(18)
    p.line_spacing = 1.5

add_footer(slide2, "Al Ahlia Contracting Group")

# --- Slide 3: Analytics Visualization ---
slide3_layout = prs.slide_layouts[5]
slide3 = prs.slides.add_slide(slide3_layout)
set_slide_background(slide3, COLOR_LIGHT_GREY)

title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
p3 = title3.text_frame.paragraphs[0]
p3.text = "📈 ADVANCED LEAVE ANALYTICS"
p3.font.name = FONT_FAMILY
p3.font.bold = True
p3.font.size = Pt(32)
p3.font.color.rgb = COLOR_DARK_BLUE

# Monthly Trend
trend_title = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
trend_title.text_frame.paragraphs[0].text = "Monthly Trend Analysis:"
trend_title.text_frame.paragraphs[0].font.bold = True
trend_title.text_frame.paragraphs[0].font.size = Pt(18)

trend_data = "Jan ▲▲▲  Feb ▲▲▲▲  Mar ▲▲  Apr ▲\nMay ▲▲▲▲▲ Jun ▲▲▲  Jul ▲▲  Aug ▲▲▲\nSep ▲▲▲▲  Oct ▲▲  Nov ▲  Dec ▲▲▲"
trend_box = slide3.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.5))
trend_frame = trend_box.text_frame
for line in trend_data.split('\n'):
    p = trend_frame.add_paragraph()
    p.text = line
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.line_spacing = 1.5

# Leave Type Distribution
type_title = slide3.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
type_title.text_frame.paragraphs[0].text = "Leave Type Distribution:"
type_title.text_frame.paragraphs[0].font.bold = True
type_title.text_frame.paragraphs[0].font.size = Pt(18)

type_data = {"Annual Leave": 45, "Sick Leave": 25, "Emergency": 15, "Other": 15}
y_pos = 4.3
for leave_type, val in type_data.items():
    slide3.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.4)).text_frame.paragraphs[0].text = f"{leave_type}:"
    bar_bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(y_pos + 0.1), Inches(4), Inches(0.2))
    bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = RGBColor(220, 220, 220); bar_bg.line.fill.background()
    bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(y_pos + 0.1), Inches(4 * (val/100)), Inches(0.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_LIGHT_BLUE; bar.line.fill.background()
    slide3.shapes.add_textbox(Inches(6.6), Inches(y_pos), Inches(1), Inches(0.4)).text_frame.paragraphs[0].text = f"{val}%"
    y_pos += 0.5
    
# Predictive Forecasting
forecast_title = slide3.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.5))
forecast_title.text_frame.paragraphs[0].text = "Predictive Forecasting:"
forecast_title.text_frame.paragraphs[0].font.bold = True
forecast_title.text_frame.paragraphs[0].font.size = Pt(18)
forecast_text = "Next 3 months: 12% increase expected\nPeak season: June & December"
forecast_box = slide3.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.8))
for line in forecast_text.split('\n'):
    forecast_box.text_frame.add_paragraph().text = line
    
add_footer(slide3, "Al Ahlia Contracting Group")

# --- Slide 4: Easy Leave Application ---
slide4_layout = prs.slide_layouts[5]
slide4 = prs.slides.add_slide(slide4_layout)
set_slide_background(slide4, COLOR_DARK_BLUE)

title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
p4 = title4.text_frame.paragraphs[0]
p4.text = "🚀 SEAMLESS LEAVE APPLICATION"
p4.font.name = FONT_FAMILY
p4.font.bold = True
p4.font.size = Pt(36)
p4.font.color.rgb = COLOR_WHITE
p4.alignment = PP_ALIGN.CENTER

# 3-Step Process
step_title = slide4.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4), Inches(0.5))
step_title.text_frame.paragraphs[0].text = "3-STEP PROCESS:"
step_title.text_frame.paragraphs[0].font.bold = True
step_title.text_frame.paragraphs[0].font.size = Pt(20)
step_title.text_frame.paragraphs[0].font.color.rgb = COLOR_GOLD

step_text = "1. 📱 CLICK \"Apply Leave\"\n2. 📅 SELECT Dates\n3. ✅ SUBMIT & Track"
step_box = slide4.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(4.5), Inches(2))
for line in step_text.split('\n'):
    p = step_box.text_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_WHITE
    p.line_spacing = 1.5

# Performance Metrics
perf_title = slide4.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(0.5))
perf_title.text_frame.paragraphs[0].text = "⏱️ PERFORMANCE METRICS:"
perf_title.text_frame.paragraphs[0].font.bold = True
perf_title.text_frame.paragraphs[0].font.size = Pt(20)
perf_title.text_frame.paragraphs[0].font.color.rgb = COLOR_GOLD

perf_text = "Before: 15 minutes per application\nAfter: 2 minutes per application\n85% Time Savings!"
perf_box = slide4.shapes.add_textbox(Inches(5.5), Inches(2.3), Inches(4.5), Inches(2))
for line in perf_text.split('\n'):
    p = perf_box.text_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_WHITE
    p.line_spacing = 1.5
    if "Savings" in line:
        p.font.bold = True
        p.font.color.rgb = COLOR_GOLD

# Mobile-First Features
mobile_title = slide4.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.5))
mobile_title.text_frame.paragraphs[0].text = "📲 MOBILE-FIRST FEATURES:"
mobile_title.text_frame.paragraphs[0].font.bold = True
mobile_title.text_frame.paragraphs[0].font.size = Pt(20)
mobile_title.text_frame.paragraphs[0].font.color.rgb = COLOR_GOLD

mobile_text = "• One-click date selection\n• Auto-balance calculation\n• Instant policy guidance\n• Push notification updates"
mobile_box = slide4.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2))
for line in mobile_text.split('\n'):
    p = mobile_box.text_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_WHITE
    p.line_spacing = 1.2

add_footer(slide4, "Al Ahlia Contracting Group")

# --- Slide 5: Application Process Flow ---
slide5_layout = prs.slide_layouts[5]
slide5 = prs.slides.add_slide(slide5_layout)
set_slide_background(slide5, COLOR_LIGHT_GREY)

title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
p5 = title5.text_frame.paragraphs[0]
p5.text = "🔄 DIGITAL WORKFLOW PROCESS"
p5.font.name = FONT_FAMILY
p5.font.bold = True
p5.font.size = Pt(32)
p5.font.color.rgb = COLOR_DARK_BLUE

# Process Flow
flow_steps = ["Employee Portal", "Smart Form", "Auto-Validation", "Manager Review", "Instant Approval", "System Update"]
x_pos, y_pos, width, height = 0.5, 2.0, 1.4, 1.0
for i, step in enumerate(flow_steps):
    shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(y_pos), Inches(width), Inches(height))
    shape.text = step
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_LIGHT_BLUE
    shape.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    if i < len(flow_steps) - 1:
        arrow = slide5.shapes.add_textbox(Inches(x_pos + width), Inches(y_pos), Inches(0.5), Inches(height))
        arrow.text_frame.paragraphs[0].text = "→"
        arrow.text_frame.paragraphs[0].font.size = Pt(36)
        arrow.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    x_pos += width + 0.5

# Automated Checks
checks_title = slide5.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(4), Inches(0.5))
checks_title.text_frame.paragraphs[0].text = "✅ AUTOMATED CHECKS:"
checks_title.text_frame.paragraphs[0].font.bold = True
checks_title.text_frame.paragraphs[0].font.size = Pt(18)

checks_text = "• Leave balance verification\n• Policy compliance check\n• Schedule conflict detection\n• Manager availability"
checks_box = slide5.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(4.5), Inches(2.5))
for line in checks_text.split('\n'):
    p = checks_box.text_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.line_spacing = 1.5

# Processing Time
time_title = slide5.shapes.add_textbox(Inches(5.5), Inches(4.0), Inches(4), Inches(0.5))
time_title.text_frame.paragraphs[0].text = "⏰ PROCESSING TIME:"
time_title.text_frame.paragraphs[0].font.bold = True
time_title.text_frame.paragraphs[0].font.size = Pt(18)

time_text = "Traditional: 3-5 days\nDigital: 2-4 hours\n94% Faster!"
time_box = slide5.shapes.add_textbox(Inches(5.5), Inches(4.5), Inches(4.5), Inches(2.5))
for line in time_text.split('\n'):
    p = time_box.text_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.line_spacing = 1.5
    if "Faster" in line:
        p.font.bold = True
        p.font.color.rgb = COLOR_LIGHT_BLUE

add_footer(slide5, "Al Ahlia Contracting Group")

# --- Slide 6: Leave Tracking & Monitoring ---
slide6_layout = prs.slide_layouts[5]
slide6 = prs.slides.add_slide(slide6_layout)
set_slide_background(slide6, COLOR_LIGHT_GREY)

title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
p6 = title6.text_frame.paragraphs[0]
p6.text = "📋 REAL-TIME TRACKING DASHBOARD"
p6.font.name = FONT_FAMILY
p6.font.bold = True
p6.font.size = Pt(32)
p6.font.color.rgb = COLOR_DARK_BLUE

# Leave Balance
balance_title = slide6.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(0.5))
balance_title.text_frame.paragraphs[0].text = "LEAVE BALANCE:"
balance_title.text_frame.paragraphs[0].font.bold = True
balance_title.text_frame.paragraphs[0].font.size = Pt(18)
# Bar
bar_bg = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.3), Inches(9), Inches(0.4))
bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = RGBColor(220, 220, 220); bar_bg.line.fill.background()
remaining_days = 18
total_days = 30
bar = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.3), Inches(9 * (remaining_days/total_days)), Inches(0.4))
bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_GOLD; bar.line.fill.background()
bar_text = slide6.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(0.4))
bar_text.text_frame.paragraphs[0].text = f"{remaining_days}/{total_days} days remaining"
bar_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
bar_text.text_frame.paragraphs[0].font.bold = True
bar_text.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_BLUE

# Status Overview & Smart Notifications
status_title = slide6.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(4), Inches(0.5))
status_title.text_frame.paragraphs[0].text = "STATUS OVERVIEW:"
status_title.text_frame.paragraphs[0].font.bold = True
status_text = "✅ Approved: 12 days\n⏳ Pending: 3 days\n📊 Used: 10 days"
status_box = slide6.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(4), Inches(1.5))
for line in status_text.split('\n'): status_box.text_frame.add_paragraph().text = line

notify_title = slide6.shapes.add_textbox(Inches(5.5), Inches(3.5), Inches(4), Inches(0.5))
notify_title.text_frame.paragraphs[0].text = "🔔 SMART NOTIFICATIONS:"
notify_title.text_frame.paragraphs[0].font.bold = True
notify_text = "• Real-time status updates\n• Low balance alerts\n• Approval notifications\n• Upcoming leave reminders"
notify_box = slide6.shapes.add_textbox(Inches(5.5), Inches(4.0), Inches(4), Inches(2))
for line in notify_text.split('\n'): notify_box.text_frame.add_paragraph().text = line

# Trend Analysis
trend_title_s6 = slide6.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(0.5))
trend_title_s6.text_frame.paragraphs[0].text = "📈 TREND ANALYSIS:"
trend_title_s6.text_frame.paragraphs[0].font.bold = True
trend_text_s6 = "• Personal leave patterns\n• Historical comparison\n• Future projections"
trend_box_s6 = slide6.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(9), Inches(1.5))
for line in trend_text_s6.split('\n'): trend_box_s6.text_frame.add_paragraph().text = line

add_footer(slide6, "Al Ahlia Contracting Group")

# --- Slide 7: System Benefits & ROI ---
slide7_layout = prs.slide_layouts[5]
slide7 = prs.slides.add_slide(slide7_layout)
set_slide_background(slide7, COLOR_DARK_BLUE)

title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
p7 = title7.text_frame.paragraphs[0]
p7.text = "💰 TRANSFORMATIVE BUSINESS IMPACT"
p7.font.name = FONT_FAMILY
p7.font.bold = True
p7.font.size = Pt(36)
p7.font.color.rgb = COLOR_WHITE
p7.alignment = PP_ALIGN.CENTER

# Operational Excellence
op_ex_title = slide7.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4), Inches(0.5))
op_ex_title.text_frame.paragraphs[0].text = "OPERATIONAL EXCELLENCE:"
op_ex_title.text_frame.paragraphs[0].font.bold = True
op_ex_title.text_frame.paragraphs[0].font.color.rgb = COLOR_GOLD
op_ex_title.text_frame.paragraphs[0].font.size = Pt(18)
op_ex_data = {"Paper Reduction": 95, "Process Speed": 85, "Cost Savings": 75, "Accuracy Rate": 98}
y_pos = 2.4
for label, val in op_ex_data.items():
    slide7.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.4)).text_frame.paragraphs[0].text = f"{label}:"
    slide7.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    bar_bg = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(y_pos + 0.1), Inches(2), Inches(0.2))
    bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = RGBColor(70, 90, 110); bar_bg.line.fill.background()
    bar = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(y_pos + 0.1), Inches(2 * (val/100)), Inches(0.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = COLOR_GOLD; bar.line.fill.background()
    slide7.shapes.add_textbox(Inches(4.6), Inches(y_pos), Inches(1), Inches(0.4)).text_frame.paragraphs[0].text = f"{val}%"
    y_pos += 0.6

# Strategic Advantages
strat_ad_title = slide7.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(0.5))
strat_ad_title.text_frame.paragraphs[0].text = "STRATEGIC ADVANTAGES:"
strat_ad_title.text_frame.paragraphs[0].font.bold = True
strat_ad_title.text_frame.paragraphs[0].font.color.rgb = COLOR_GOLD
strat_ad_title.text_frame.paragraphs[0].font.size = Pt(18)
strat_ad_text = "• HR workload reduced by 70%\n• Employee satisfaction: 92%\n• Compliance accuracy: 100%\n• Mobile adoption: 95%"
strat_ad_box = slide7.shapes.add_textbox(Inches(5.5), Inches(2.4), Inches(4), Inches(4))
for line in strat_ad_text.split('\n'):
    p = strat_ad_box.text_frame.add_paragraph()
    p.text = line
    p.line_spacing = 1.5
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_WHITE

add_footer(slide7, "Al Ahlia Contracting Group")

# --- Slide 8: Quantifiable Results ---
slide8_layout = prs.slide_layouts[5]
slide8 = prs.slides.add_slide(slide8_layout)
set_slide_background(slide8, COLOR_LIGHT_GREY)

title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
p8 = title8.text_frame.paragraphs[0]
p8.text = "📈 MEASURABLE PERFORMANCE"
p8.font.name = FONT_FAMILY
p8.font.bold = True
p8.font.size = Pt(32)
p8.font.color.rgb = COLOR_DARK_BLUE

# Before vs After
bva_title = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
bva_title.text_frame.paragraphs[0].text = "BEFORE vs AFTER:"
bva_title.text_frame.paragraphs[0].font.bold = True
bva_title.text_frame.paragraphs[0].font.size = Pt(18)
bva_text = "Processing Time: 3 days → 2 hours (94% faster)\nEmployee Satisfaction: 65% → 92% (+27 pts)\nHR Workload: 40 hrs/wk → 12 hrs/wk (70% reduction)\nError Rate: 15% → 2% (87% improvement)\nCost/Transaction: $25 → $8 (68% savings)"
bva_box = slide8.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(2.5))
for line in bva_text.split('\n'):
    p = bva_box.text_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.line_spacing = 1.5

# Key Performance Indicators
kpi_title = slide8.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.5))
kpi_title.text_frame.paragraphs[0].text = "KEY PERFORMANCE INDICATORS:"
kpi_title.text_frame.paragraphs[0].font.bold = True
kpi_title.text_frame.paragraphs[0].font.size = Pt(18)
kpi_text = "⏱️ Time Efficiency: 94% faster\n💸 Cost Reduction: 68% lower\n😊 Employee Satisfaction: 92% rating\n📈 Adoption Rate: 98% engagement"
kpi_box = slide8.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(2.0))
for line in kpi_text.split('\n'):
    p = kpi_box.text_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.line_spacing = 1.5
    
add_footer(slide8, "Al Ahlia Contracting Group")

# --- Slide 9: Visual Analytics Dashboard ---
slide9_layout = prs.slide_layouts[5]
slide9 = prs.slides.add_slide(slide9_layout)
set_slide_background(slide9, COLOR_LIGHT_GREY)

title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
p9 = title9.text_frame.paragraphs[0]
p9.text = "🎨 LIVE ANALYTICS DEMO"
p9.font.name = FONT_FAMILY
p9.font.bold = True
p9.font.size = Pt(32)
p9.font.color.rgb = COLOR_DARK_BLUE

# Placeholder
placeholder = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
placeholder.fill.solid(); placeholder.fill.fore_color.rgb = RGBColor(220, 220, 220)
placeholder.text = "[INSERT ACTUAL SYSTEM SCREENSHOT HERE]"
placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Real-time Metrics
metrics_title = slide9.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(3), Inches(0.5))
metrics_title.text_frame.paragraphs[0].text = "REAL-TIME METRICS:"
metrics_title.text_frame.paragraphs[0].font.bold = True
metrics_text = "• Current leave applications: 12\n• Pending approvals: 3\n• This month's utilization: 42%\n• Department availability: 78%"
metrics_box = slide9.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(3), Inches(2.0))
for line in metrics_text.split('\n'): metrics_box.text_frame.add_paragraph().text = line

# Predictive Insights
insights_title = slide9.shapes.add_textbox(Inches(6.8), Inches(4.5), Inches(3), Inches(0.5))
insights_title.text_frame.paragraphs[0].text = "PREDICTIVE INSIGHTS:"
insights_title.text_frame.paragraphs[0].font.bold = True
insights_text = "• Next week: 8 expected apps\n• Peak day: Friday approvals\n• Busiest department: Sales"
insights_box = slide9.shapes.add_textbox(Inches(6.8), Inches(5.0), Inches(3), Inches(2.0))
for line in insights_text.split('\n'): insights_box.text_frame.add_paragraph().text = line

add_footer(slide9, "Al Ahlia Contracting Group")

# --- Slide 10: Future Roadmap ---
slide10_layout = prs.slide_layouts[5]
slide10 = prs.slides.add_slide(slide10_layout)
set_slide_background(slide10, COLOR_DARK_BLUE)

title10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
p10 = title10.text_frame.paragraphs[0]
p10.text = "🗺️ CONTINUOUS INNOVATION"
p10.font.name = FONT_FAMILY
p10.font.bold = True
p10.font.size = Pt(36)
p10.font.color.rgb = COLOR_WHITE
p10.alignment = PP_ALIGN.CENTER

# Roadmap sections
roadmap_data = {
    "NEXT 6 MONTHS: 🤖 AI-Powered Predictions": "• Smart leave recommendations\n• Automated conflict resolution\n• Predictive analytics",
    "🔗 Integration Expansion": "• Payroll system sync\n• Performance management\n• Enhanced mobile app",
    "🌍 Global Compliance": "• Multi-country policies\n• Local labor law updates\n• International calendars"
}
y_pos = 1.8
x_pos = 0.5
width = 3.0
for i, (title, text) in enumerate(roadmap_data.items()):
    title_box = slide10.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(width), Inches(0.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.color.rgb = COLOR_GOLD
    p.font.size = Pt(16)
    
    text_box = slide10.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.5), Inches(width), Inches(2.5))
    for line in text.split('\n'):
        p_text = text_box.text_frame.add_paragraph()
        p_text.text = line
        p_text.font.color.rgb = COLOR_WHITE
        p_text.font.size = Pt(14)
    x_pos += width + 0.2

# Vision 2025
vision_box = slide10.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(1.0))
vision_box.text = 'VISION 2025:\n"Most intelligent leave platform in the region"'
p_vision_title = vision_box.text_frame.paragraphs[0]
p_vision_title.font.bold = True
p_vision_title.font.color.rgb = COLOR_GOLD
p_vision_title.font.size = Pt(18)
p_vision_title.alignment = PP_ALIGN.CENTER
p_vision_quote = vision_box.text_frame.paragraphs[1]
p_vision_quote.font.italic = True
p_vision_quote.font.color.rgb = COLOR_WHITE
p_vision_quote.font.size = Pt(18)
p_vision_quote.alignment = PP_ALIGN.CENTER

add_footer(slide10, "Al Ahlia Contracting Group")

# --- Slide 11: Thank You ---
slide11_layout = prs.slide_layouts[5]
slide11 = prs.slides.add_slide(slide11_layout)
set_slide_background(slide11, COLOR_DARK_BLUE)

thank_you_shape = slide11.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
thank_you_para = thank_you_shape.text_frame.paragraphs[0]
thank_you_para.text = "🙏 THANK YOU"
thank_you_para.font.name = FONT_FAMILY
thank_you_para.font.size = Pt(60)
thank_you_para.font.bold = True
thank_you_para.font.color.rgb = COLOR_WHITE
thank_you_para.alignment = PP_ALIGN.CENTER

discussion_shape = slide11.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(1.0))
discussion_para = discussion_shape.text_frame.paragraphs[0]
discussion_para.text = "Questions & Discussion"
discussion_para.font.name = FONT_FAMILY
discussion_para.font.size = Pt(28)
discussion_para.font.color.rgb = COLOR_GOLD
discussion_para.alignment = PP_ALIGN.CENTER

company_shape = slide11.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.0))
company_para = company_shape.text_frame.paragraphs[0]
company_para.text = "Al Ahlia Contracting Group"
company_para.font.name = FONT_FAMILY
company_para.font.size = Pt(24)
company_para.font.color.rgb = COLOR_WHITE
company_para.alignment = PP_ALIGN.CENTER


# --- Save Presentation ---
file_path = "HR_Leave_Management_System.pptx"
prs.save(file_path)

print(f"Presentation saved to {file_path}")
